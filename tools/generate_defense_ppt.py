from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(r"D:\falldetection")
OUT_PATH = ROOT / "闫家乐-毕业答辩PPT-生成版.pptx"

TITLE = "跌倒检测与事件留证系统的设计与实现"
AUTHOR = "闫家乐"
STUDENT_ID = "2406249044"
MAJOR = "计算机科学与技术"
COLLEGE = "计算科学与人工智能学院"
ADVISOR = "许晨航"
SUBMIT_DATE = "2026年4月"

SCREEN_DIR = ROOT / "thesis_assets" / "screenshots"
ASSET_DIR = ROOT / "thesis_assets"


COLORS = {
    "navy": RGBColor(15, 23, 42),
    "slate": RGBColor(30, 41, 59),
    "purple": RGBColor(109, 40, 217),
    "cyan": RGBColor(6, 182, 212),
    "green": RGBColor(16, 185, 129),
    "red": RGBColor(239, 68, 68),
    "amber": RGBColor(245, 158, 11),
    "text": RGBColor(30, 41, 59),
    "muted": RGBColor(100, 116, 139),
    "light": RGBColor(248, 250, 252),
    "line": RGBColor(226, 232, 240),
    "white": RGBColor(255, 255, 255),
}


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_bg(slide, color=None):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color or COLORS["light"]


def add_top_band(slide, section_no, title, subtitle=""):
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.82))
    band.fill.solid()
    band.fill.fore_color.rgb = COLORS["navy"]
    band.line.fill.background()

    num_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.38), Inches(0.18), Inches(0.72), Inches(0.42))
    num_box.fill.solid()
    num_box.fill.fore_color.rgb = COLORS["purple"]
    num_box.line.fill.background()
    tf = num_box.text_frame
    tf.text = section_no
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.name = "Microsoft YaHei"
    run.font.bold = True
    run.font.size = Pt(20)
    run.font.color.rgb = COLORS["white"]

    title_box = slide.shapes.add_textbox(Inches(1.3), Inches(0.12), Inches(5.6), Inches(0.34))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Microsoft YaHei"
    r.font.bold = True
    r.font.size = Pt(24)
    r.font.color.rgb = COLORS["white"]

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(1.32), Inches(0.42), Inches(7.5), Inches(0.18))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = subtitle
        r.font.name = "Microsoft YaHei"
        r.font.size = Pt(10.5)
        r.font.color.rgb = RGBColor(203, 213, 225)


def add_footer(slide, text):
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.4), Inches(7.1), Inches(12.55), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["line"]
    line.line.fill.background()
    box = slide.shapes.add_textbox(Inches(10.8), Inches(7.12), Inches(2.0), Inches(0.2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = text
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(9)
    r.font.color.rgb = COLORS["muted"]


def add_title_text(slide, text, left, top, width, height, size=24, color=None, bold=True):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color or COLORS["text"]
    return box


def add_bullets(slide, items, left, top, width, height, font_size=17, color=None, bullet_color=None, line_gap=1.15):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.line_spacing = line_gap
        p.space_after = Pt(5)
        p.bullet = True
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color or COLORS["text"]
        if bullet_color:
            for run in p.runs:
                run.font.color.rgb = color or COLORS["text"]
    return box


def add_card(slide, left, top, width, height, title, body_lines, accent):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = COLORS["white"]
    card.line.color.rgb = COLORS["line"]
    card.line.width = Pt(1.0)
    stripe = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, Inches(0.08), height)
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = accent
    stripe.line.fill.background()
    add_title_text(slide, title, left + Inches(0.18), top + Inches(0.15), width - Inches(0.28), Inches(0.28), size=18)
    add_bullets(slide, body_lines, left + Inches(0.14), top + Inches(0.52), width - Inches(0.24), height - Inches(0.62), font_size=13.5)


def add_tag(slide, text, left, top, width, fill):
    tag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, Inches(0.32))
    tag.fill.solid()
    tag.fill.fore_color.rgb = fill
    tag.line.fill.background()
    tf = tag.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(10.5)
    run.font.bold = True
    run.font.color.rgb = COLORS["white"]


def add_picture(slide, path, left, top, width=None, height=None):
    if not Path(path).exists():
        return None
    return slide.shapes.add_picture(str(path), left, top, width=width, height=height)


def add_callout(slide, left, top, width, text, fill):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, Inches(0.48))
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = COLORS["white"]


def cover_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, COLORS["light"])

    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.18))
    band.fill.solid()
    band.fill.fore_color.rgb = COLORS["navy"]
    band.line.fill.background()

    add_title_text(slide, "本科生毕业论文（设计）答辩", Inches(0.6), Inches(0.28), Inches(5.2), Inches(0.4), size=22, color=COLORS["white"])
    add_title_text(slide, TITLE, Inches(0.72), Inches(1.45), Inches(7.6), Inches(0.95), size=29)

    sub = slide.shapes.add_textbox(Inches(0.75), Inches(2.35), Inches(5.6), Inches(0.42))
    tf = sub.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "开场：各位老师好，我的题目是《跌倒检测与事件留证系统的设计与实现》，我是闫家乐，我的指导老师是许晨航老师。"
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(13)
    r.font.color.rgb = COLORS["muted"]

    info = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.72), Inches(3.0), Inches(5.4), Inches(2.2))
    info.fill.solid()
    info.fill.fore_color.rgb = COLORS["white"]
    info.line.color.rgb = COLORS["line"]

    entries = [
        ("学院", COLLEGE),
        ("专业", MAJOR),
        ("姓名", AUTHOR),
        ("学号", STUDENT_ID),
        ("指导教师", ADVISOR),
        ("提交时间", SUBMIT_DATE),
    ]
    top = 3.22
    for label, value in entries:
        add_title_text(slide, label, Inches(1.02), Inches(top), Inches(1.0), Inches(0.24), size=13, color=COLORS["muted"], bold=False)
        add_title_text(slide, value, Inches(2.05), Inches(top - 0.02), Inches(3.65), Inches(0.28), size=15)
        top += 0.31

    add_tag(slide, "答辩重点：需求 - 设计 - 实现", Inches(0.75), Inches(5.48), Inches(2.45), COLORS["purple"])
    add_tag(slide, "重点模块：实时检测 / 事件留证", Inches(3.36), Inches(5.48), Inches(2.85), COLORS["cyan"])

    add_picture(slide, SCREEN_DIR / "detect-5173.png", Inches(7.1), Inches(1.44), width=Inches(5.6), height=Inches(3.24))
    add_picture(slide, ASSET_DIR / "sequence_diagram_alarm_flow.png", Inches(7.45), Inches(4.9), width=Inches(5.0), height=Inches(1.85))
    add_footer(slide, "封面")


def slide_requirements_1():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_top_band(slide, "01", "需求分析", "打算做什么：先解决养老监护场景中的实时识别、及时告警和证据留存问题")

    add_card(
        slide, Inches(0.55), Inches(1.2), Inches(4.0), Inches(2.15), "建设目标",
        ["对摄像头或本地视频中的跌倒行为进行实时识别", "识别后立即告警，并记录事件前后的视频证据", "支持管理员介入处理，形成完整业务闭环"],
        COLORS["purple"],
    )
    add_card(
        slide, Inches(4.72), Inches(1.2), Inches(4.0), Inches(2.15), "用户端需求",
        ["登录后进入检测页面，切换摄像头或上传视频", "查看告警提示、历史录像和个人资料", "配置紧急联系人，接收处理反馈"],
        COLORS["cyan"],
    )
    add_card(
        slide, Inches(8.9), Inches(1.2), Inches(3.85), Inches(2.15), "管理端需求",
        ["查看告警记录与处理状态", "配置通知方式、时段和联系人", "管理统计分析、日志中心、模型训练与系统参数"],
        COLORS["green"],
    )

    add_card(
        slide, Inches(0.55), Inches(3.7), Inches(6.15), Inches(2.75), "关键非功能需求",
        ["实时性：检测后要快速弹出预警，减少延迟", "可用性：支持摄像头和本地视频两类演示入口", "安全性：区分管理员与普通用户的功能边界", "可维护性：前后端分离，模块化扩展业务功能"],
        COLORS["amber"],
    )

    add_picture(slide, ASSET_DIR / "use_case_diagram.png", Inches(7.05), Inches(3.68), width=Inches(5.35), height=Inches(2.72))
    add_tag(slide, "需求结论：系统不能只识别，还必须能告警、留证、回放、处置", Inches(7.0), Inches(6.48), Inches(5.45), COLORS["navy"])
    add_footer(slide, "01 / 需求分析")


def slide_requirements_2():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_top_band(slide, "02", "需求分析", "围绕真实业务流程拆分：检测识别 -> 告警通知 -> 证据留存 -> 人工处置")

    add_title_text(slide, "核心业务链路", Inches(0.62), Inches(1.18), Inches(2.2), Inches(0.3), size=20)
    steps = [
        ("1 视频输入", "支持摄像头与本地文件，方便演示和部署"),
        ("2 模型推理", "YOLO 模型识别人体姿态与跌倒状态"),
        ("3 触发报警", "前端弹窗，后端保存报警记录"),
        ("4 事件留证", "保存跌倒前后视频片段和关键帧截图"),
        ("5 人工处理", "管理员登记工单结果，完善处置闭环"),
    ]
    y = 1.72
    accents = [COLORS["purple"], COLORS["cyan"], COLORS["red"], COLORS["green"], COLORS["amber"]]
    for idx, (title, desc) in enumerate(steps):
        add_card(slide, Inches(0.68), Inches(y), Inches(5.62), Inches(0.78), title, [desc], accents[idx])
        y += 0.92

    add_title_text(slide, "答辩重点为什么放在这两块", Inches(6.75), Inches(1.18), Inches(3.8), Inches(0.3), size=20)
    add_card(
        slide, Inches(6.82), Inches(1.72), Inches(5.7), Inches(1.55), "重点一：实时检测与事件留证",
        ["工作量最大，涉及模型推理、缓冲区、视频写入、截图提取和回放索引"], COLORS["purple"]
    )
    add_card(
        slide, Inches(6.82), Inches(3.53), Inches(5.7), Inches(1.55), "重点二：报警管理与工单闭环",
        ["更能体现系统价值，把“识别结果”变成“可处理的业务事件”"], COLORS["cyan"]
    )
    add_card(
        slide, Inches(6.82), Inches(5.34), Inches(5.7), Inches(1.05), "本页结论",
        ["所以后面的设计与实现，主要围绕这两个模块展开说明"], COLORS["navy"]
    )
    add_footer(slide, "02 / 需求分析")


def slide_architecture():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_top_band(slide, "03", "系统设计", "怎么做：采用前后端分离架构，并围绕核心业务划分功能模块")

    stage_w = Inches(2.05)
    stage_h = Inches(1.12)
    titles = ["前端展示层", "后端服务层", "检测与存储层", "数据与证据层"]
    notes = [
        "Vue 3 页面交互\n检测、报警、回放、统计",
        "Flask API 与权限控制\n接口组织、业务编排",
        "YOLO + OpenCV + 缓冲区\n实时识别与视频处理",
        "MongoDB / GridFS\n记录、截图、留证视频",
    ]
    x_positions = [0.75, 3.25, 5.75, 8.25]
    for idx in range(4):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_positions[idx]), Inches(1.65), stage_w, stage_h)
        card.fill.solid()
        card.fill.fore_color.rgb = COLORS["white"]
        card.line.color.rgb = COLORS["line"]
        add_title_text(slide, titles[idx], Inches(x_positions[idx] + 0.12), Inches(1.82), Inches(1.7), Inches(0.25), size=16)
        text = notes[idx].split("\n")
        add_bullets(slide, text, Inches(x_positions[idx] + 0.08), Inches(2.1), Inches(1.85), Inches(0.56), font_size=10.5, line_gap=1.0)
        if idx < 3:
            arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, Inches(x_positions[idx] + 2.1), Inches(2.02), Inches(0.5), Inches(0.32))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = COLORS["cyan"]
            arrow.line.fill.background()

    add_title_text(slide, "功能模块划分", Inches(0.62), Inches(3.18), Inches(2.4), Inches(0.3), size=20)
    modules = [
        ("登录认证", COLORS["purple"]),
        ("实时检测", COLORS["cyan"]),
        ("报警管理", COLORS["red"]),
        ("历史回放", COLORS["green"]),
        ("统计/日志/训练", COLORS["amber"]),
        ("系统配置", COLORS["navy"]),
    ]
    start_x = 0.8
    start_y = 3.75
    for idx, (name, color) in enumerate(modules):
        col = idx % 3
        row = idx // 3
        left = Inches(start_x + col * 2.15)
        top = Inches(start_y + row * 1.0)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.9), Inches(0.72))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        tf = shape.text_frame
        tf.text = name
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.name = "Microsoft YaHei"
        run.font.size = Pt(15)
        run.font.bold = True
        run.font.color.rgb = COLORS["white"]

    add_picture(slide, ASSET_DIR / "class_diagram.png", Inches(7.55), Inches(3.2), width=Inches(4.8), height=Inches(3.0))
    add_tag(slide, "设计思路：先保证检测链路跑通，再补齐管理功能与证据闭环", Inches(0.78), Inches(6.52), Inches(6.1), COLORS["navy"])
    add_footer(slide, "03 / 系统设计")


def slide_module_detection():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_top_band(slide, "04", "重点模块设计", "模块一：实时检测与事件留证，是系统最核心、最有工作量的部分")

    add_picture(slide, ASSET_DIR / "sequence_diagram_alarm_flow.png", Inches(7.25), Inches(1.28), width=Inches(5.1), height=Inches(4.45))
    add_title_text(slide, "设计要点", Inches(0.62), Inches(1.2), Inches(2.1), Inches(0.3), size=20)
    add_card(
        slide, Inches(0.68), Inches(1.7), Inches(5.75), Inches(1.2), "1 检测入口统一",
        ["前端支持摄像头与本地视频两种输入方式，便于答辩演示与后续部署"], COLORS["purple"]
    )
    add_card(
        slide, Inches(0.68), Inches(3.02), Inches(5.75), Inches(1.35), "2 留证机制不是单纯录像",
        ["后端维护跌倒前循环缓冲区和跌倒后补录区，确保保留下跌前后关键片段"], COLORS["cyan"]
    )
    add_card(
        slide, Inches(0.68), Inches(4.49), Inches(5.75), Inches(1.35), "3 数据落库形成回放基础",
        ["视频写入 GridFS，history 集合保存索引，alarm_snapshots 保存关键帧截图"], COLORS["green"]
    )

    add_tag(slide, "难点：既要实时识别，又要把事件发生前后的现场完整保留下来", Inches(0.72), Inches(6.2), Inches(5.8), COLORS["red"])
    add_footer(slide, "04 / 重点模块设计")


def slide_module_alarm():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_top_band(slide, "05", "重点模块设计", "模块二：报警管理与工单闭环，把算法结果转成可处理的业务记录")

    add_picture(slide, SCREEN_DIR / "alarm-5173.png", Inches(7.05), Inches(1.22), width=Inches(5.48), height=Inches(4.95))
    add_title_text(slide, "模块价值", Inches(0.62), Inches(1.2), Inches(2.0), Inches(0.3), size=20)
    add_card(
        slide, Inches(0.7), Inches(1.7), Inches(5.72), Inches(1.12), "1 告警配置集中管理",
        ["管理员可以配置声音、浏览器、邮件、短信以及通知时段和联系人"], COLORS["purple"]
    )
    add_card(
        slide, Inches(0.7), Inches(2.95), Inches(5.72), Inches(1.12), "2 历史告警可筛选、可追踪",
        ["按状态和日期筛选记录，查看待处理、已处理及详情信息"], COLORS["cyan"]
    )
    add_card(
        slide, Inches(0.7), Inches(4.2), Inches(5.72), Inches(1.12), "3 工单机制补齐人工处置",
        ["管理员提交真实跌倒、误报或继续观察等结果，形成处理闭环"], COLORS["amber"]
    )
    add_card(
        slide, Inches(0.7), Inches(5.45), Inches(5.72), Inches(0.82), "创新点",
        ["系统不只提示告警，还保留处置结果、通知日志与审核痕迹"], COLORS["navy"]
    )
    add_footer(slide, "05 / 重点模块设计")


def slide_ui_overview():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_top_band(slide, "06", "系统实现", "3 到 4 页展示系统界面：这一页先给出整体功能视图")

    pics = [
        ("实时检测", SCREEN_DIR / "detect-5173.png"),
        ("报警管理", SCREEN_DIR / "alarm-5173.png"),
        ("历史回放", SCREEN_DIR / "replay-5173.png"),
        ("统计分析", SCREEN_DIR / "statistics-5173.png"),
    ]
    positions = [
        (0.62, 1.25),
        (6.72, 1.25),
        (0.62, 4.0),
        (6.72, 4.0),
    ]
    for (label, path), (x, y) in zip(pics, positions):
        frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(5.35), Inches(2.18))
        frame.fill.solid()
        frame.fill.fore_color.rgb = COLORS["white"]
        frame.line.color.rgb = COLORS["line"]
        add_picture(slide, path, Inches(x + 0.08), Inches(y + 0.08), width=Inches(5.19), height=Inches(1.76))
        add_tag(slide, label, Inches(x + 0.12), Inches(y + 1.88), Inches(1.15), COLORS["navy"])

    add_tag(slide, "界面实现覆盖检测、告警、回放、统计等主要业务页面", Inches(3.7), Inches(6.55), Inches(5.85), COLORS["purple"])
    add_footer(slide, "06 / 系统实现")


def slide_ui_detail_detection():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_top_band(slide, "07", "系统实现", "重点界面一：实时检测页，体现视频输入、参数配置、结果显示与 AI 分析")

    add_picture(slide, SCREEN_DIR / "detect-5173.png", Inches(0.7), Inches(1.28), width=Inches(8.15), height=Inches(5.6))
    add_card(
        slide, Inches(9.02), Inches(1.44), Inches(3.45), Inches(1.2), "界面实现点 1",
        ["左侧统一放置任务、模型、输入源和阈值参数，演示时操作路径清晰"], COLORS["purple"]
    )
    add_card(
        slide, Inches(9.02), Inches(2.87), Inches(3.45), Inches(1.2), "界面实现点 2",
        ["中间区域展示视频流，检测开始后可直接看到识别状态变化"], COLORS["cyan"]
    )
    add_card(
        slide, Inches(9.02), Inches(4.3), Inches(3.45), Inches(1.2), "界面实现点 3",
        ["下方表格同步输出置信度和空间信息，便于调参与结果核查"], COLORS["green"]
    )
    add_callout(slide, Inches(1.08), Inches(1.58), Inches(1.45), "输入源与参数", COLORS["purple"])
    add_callout(slide, Inches(5.0), Inches(1.76), Inches(1.3), "视频主视图", COLORS["cyan"])
    add_callout(slide, Inches(4.88), Inches(5.82), Inches(1.4), "检测结果表", COLORS["green"])
    add_footer(slide, "07 / 系统实现")


def slide_ui_detail_alarm_replay():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_top_band(slide, "08", "系统实现", "重点界面二：报警管理与历史回放，体现留证结果可追踪、可复核、可回看")

    add_picture(slide, SCREEN_DIR / "alarm-5173.png", Inches(0.65), Inches(1.45), width=Inches(6.0), height=Inches(4.0))
    add_picture(slide, SCREEN_DIR / "replay-5173.png", Inches(6.9), Inches(1.45), width=Inches(5.7), height=4.0)
    add_card(
        slide, Inches(0.74), Inches(5.7), Inches(5.7), Inches(0.82), "报警页实现说明",
        ["左侧负责通知配置，右侧负责历史记录、状态筛选、详情与工单处理"], COLORS["red"]
    )
    add_card(
        slide, Inches(6.96), Inches(5.7), Inches(5.38), Inches(0.82), "回放页实现说明",
        ["通过历史记录索引快速定位留证视频，支撑事件复盘与答辩演示"], COLORS["navy"]
    )
    add_callout(slide, Inches(1.0), Inches(1.76), Inches(1.25), "告警配置", COLORS["purple"])
    add_callout(slide, Inches(4.62), Inches(1.78), Inches(1.25), "记录处理", COLORS["cyan"])
    add_callout(slide, Inches(7.2), Inches(1.76), Inches(1.25), "录像列表", COLORS["green"])
    add_callout(slide, Inches(10.1), Inches(2.0), Inches(1.25), "回放区域", COLORS["amber"])
    add_footer(slide, "08 / 系统实现")


def slide_summary():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_top_band(slide, "09", "总结与答辩结束", "不求面面俱到，重点回答系统做成了什么、难点在哪里、价值体现在哪里")

    add_card(
        slide, Inches(0.72), Inches(1.45), Inches(3.7), Inches(2.6), "本课题完成情况",
        ["完成实时检测、报警管理、历史回放、统计、训练与配置等主要模块", "实现从识别到留证再到人工处置的完整业务链路"], COLORS["purple"]
    )
    add_card(
        slide, Inches(4.78), Inches(1.45), Inches(3.7), Inches(2.6), "重点与创新",
        ["把跌倒识别与证据留存结合起来，不停留在单纯算法演示", "引入角色权限、工单处理和通知日志，增强工程落地性"], COLORS["cyan"]
    )
    add_card(
        slide, Inches(8.82), Inches(1.45), Inches(3.7), Inches(2.6), "不足与改进",
        ["后续可继续优化模型精度、长时稳定性与多设备并发能力", "进一步规范认证安全和通知链路的部署方式"], COLORS["amber"]
    )

    center = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.45), Inches(4.55), Inches(10.25), Inches(1.08))
    center.fill.solid()
    center.fill.fore_color.rgb = COLORS["navy"]
    center.line.fill.background()
    tf = center.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "谢谢各位老师，请批评指正"
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = COLORS["white"]

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = "答辩陈述建议围绕“需求 -> 设计 -> 两个重点模块 -> 界面实现 -> 总结”展开"
    r2.font.name = "Microsoft YaHei"
    r2.font.size = Pt(11.5)
    r2.font.color.rgb = RGBColor(203, 213, 225)
    add_footer(slide, "09 / 总结")


cover_slide()
slide_requirements_1()
slide_requirements_2()
slide_architecture()
slide_module_detection()
slide_module_alarm()
slide_ui_overview()
slide_ui_detail_detection()
slide_ui_detail_alarm_replay()
slide_summary()

prs.save(str(OUT_PATH))
print(f"Saved to: {OUT_PATH}")
