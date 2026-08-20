from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
ASSETS = ROOT / "thesis_assets"
SHOTS = ASSETS / "screenshots"
OUT = ROOT / "跌倒检测与事件留证系统-答辩PPT-第一版.pptx"

COLORS = {
    "bg": RGBColor(247, 247, 245),
    "navy": RGBColor(47, 72, 88),
    "teal": RGBColor(58, 125, 124),
    "accent": RGBColor(224, 122, 95),
    "gold": RGBColor(242, 204, 143),
    "text": RGBColor(31, 41, 51),
    "muted": RGBColor(82, 96, 109),
    "line": RGBColor(217, 226, 236),
    "white": RGBColor(255, 255, 255),
    "soft": RGBColor(238, 244, 243),
}


def set_bg(slide, color="bg"):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS[color]


def add_rect(slide, x, y, w, h, fill, line=None, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = COLORS[fill]
    if line:
        s.line.color.rgb = COLORS[line]
        s.line.width = Pt(0.8)
    else:
        s.line.fill.background()
    return s


def add_text(slide, text, x, y, w, h, size=18, color="text", bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = COLORS[color]
    return box


def add_bullets(slide, items, x, y, w, h, size=17, color="text"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(8)
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(size)
        p.font.color.rgb = COLORS[color]
    return box


def add_picture(slide, path, x, y, w, h):
    if not Path(path).exists():
        return None
    return slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))


def title_bar(slide, title, subtitle, n):
    add_rect(slide, 0, 0, 13.333, 0.58, "navy", radius=False)
    add_text(slide, title, 0.35, 0.08, 7.6, 0.36, 24, "white", True)
    add_text(slide, subtitle, 8.4, 0.12, 3.8, 0.3, 11, "gold", False, PP_ALIGN.RIGHT)
    add_text(slide, f"{n:02d}", 12.35, 0.08, 0.5, 0.32, 20, "gold", True, PP_ALIGN.CENTER)


def caption(slide, text, x, y, w):
    add_text(slide, text, x, y, w, 0.22, 10, "muted", False, PP_ALIGN.CENTER)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    s = prs.slides.add_slide(blank)
    set_bg(s)
    add_rect(s, 0.72, 1.06, 5.3, 4.55, "navy", radius=False)
    add_rect(s, 1.02, 1.36, 5.2, 4.45, "teal")
    add_text(s, "跌倒检测与事件留证系统的设计与实现", 0.9, 1.62, 4.9, 1.7, 28, "white", True)
    add_text(s, "本科毕业设计答辩汇报", 0.92, 3.45, 3.6, 0.35, 18, "gold")
    add_text(s, "基于项目实现与论文定稿内容整理", 0.92, 3.85, 4.2, 0.3, 13, "white")
    add_rect(s, 0.92, 4.5, 1.8, 0.05, "accent", radius=False)
    add_text(s, "研究方向：智能养老 / 计算机视觉 / Web系统设计", 0.9, 6.35, 5.4, 0.3, 14, "muted")
    add_picture(s, SHOTS / "detect-5173.png", 6.7, 1.15, 5.3, 4.2)
    caption(s, "系统实时检测界面", 6.7, 5.48, 5.3)

    s = prs.slides.add_slide(blank)
    set_bg(s)
    title_bar(s, "1. 研究背景及意义", "老龄化监护需求推动非接触式跌倒检测研究", 2)
    add_rect(s, 0.42, 0.86, 5.55, 5.8, "white", "line")
    add_text(s, "课题背景", 0.65, 1.08, 1.8, 0.35, 22, "teal", True)
    add_bullets(s, ["人口老龄化加深，养老场景对安全监护提出更高要求。", "跌倒事件突发性强、危害高，对救援时效要求高。", "穿戴式方案依赖主动佩戴，长期监护中容易漏戴和误报。"], 0.65, 1.52, 5.0, 1.75)
    add_text(s, "研究意义", 0.65, 3.68, 1.8, 0.35, 22, "teal", True)
    add_bullets(s, ["基于视觉的非接触式识别更适合连续看护场景。", "将跌倒识别、报警、留证、回放和人工处理串成闭环。", "为智慧养老系统工程化落地提供参考。"], 0.65, 4.12, 5.0, 1.65)
    add_picture(s, ASSETS / "use_case_diagram.png", 6.7, 1.28, 5.0, 3.9)
    caption(s, "系统参与者与核心业务关系图", 6.7, 5.33, 5.0)
    add_rect(s, 7.1, 5.8, 4.45, 0.72, "soft", "line")
    add_text(s, "目标：构建面向养老监护场景的跌倒检测与事件留证系统", 7.28, 5.98, 4.1, 0.36, 17, "navy", True, PP_ALIGN.CENTER)

    s = prs.slides.add_slide(blank)
    set_bg(s)
    title_bar(s, "2. 研究方法及过程", "技术路线与总体架构", 3)
    add_rect(s, 0.42, 0.9, 11.9, 5.7, "white", "line")
    add_text(s, "总体技术路线", 0.65, 1.14, 2.5, 0.35, 23, "teal", True)
    add_bullets(s, ["核心识别：YOLO 目标检测模型 + OpenCV 视频流处理。", "系统架构：Vue 3 前端 + Flask 后端 + MongoDB / GridFS 存储。", "运行方式：支持摄像头与本地视频输入，覆盖检测、报警、留证、回放、统计与训练。"], 0.65, 1.62, 5.65, 1.85)
    add_text(s, "关键研究过程", 0.65, 3.78, 2.5, 0.35, 23, "teal", True)
    add_bullets(s, ["需求分析：明确实时性、可用性、安全性与可维护性目标。", "总体设计：划分检测、报警、回放、用户中心、日志与训练模块。", "详细实现：完成接口、权限、数据库结构和留证链路。"], 0.65, 4.25, 5.65, 1.65)
    add_picture(s, ASSETS / "er_diagram.png", 6.85, 1.25, 5.1, 4.2)
    caption(s, "数据库与业务对象关系示意", 6.85, 5.62, 5.1)

    s = prs.slides.add_slide(blank)
    set_bg(s)
    title_bar(s, "2. 研究方法及过程", "功能模块设计", 4)
    add_rect(s, 0.42, 0.95, 3.8, 5.5, "white", "line")
    add_text(s, "系统功能模块", 0.65, 1.18, 2.5, 0.35, 23, "teal", True)
    add_bullets(s, ["实时检测：视频输入、参数调整、识别结果展示、报警触发。", "报警管理：报警记录、联系人配置、通知时段、工单处理。", "历史回放：视频片段留存、关键帧截图、记录删除与回放。", "管理模块：用户资料、日志中心、统计分析、模型训练、系统配置。"], 0.65, 1.68, 3.25, 3.05, 16)
    add_rect(s, 0.65, 5.08, 3.25, 1.05, "soft", "line")
    add_text(s, "模块设计强调形成完整的事件处理与业务管理闭环。", 0.82, 5.33, 2.9, 0.45, 16, "navy", True, PP_ALIGN.CENTER)
    add_picture(s, SHOTS / "detect-5173.png", 4.55, 1.2, 3.35, 2.2)
    add_picture(s, SHOTS / "alarm-5173.png", 8.5, 1.2, 3.35, 2.2)
    add_picture(s, SHOTS / "replay-5173.png", 4.55, 4.02, 3.35, 2.05)
    add_picture(s, SHOTS / "statistics-5173.png", 8.5, 4.02, 3.35, 2.05)

    s = prs.slides.add_slide(blank)
    set_bg(s)
    title_bar(s, "2. 研究方法及过程", "关键业务流程", 5)
    add_picture(s, ASSETS / "sequence_diagram_alarm_flow.png", 0.7, 1.22, 5.8, 4.58)
    caption(s, "报警处理时序图", 0.7, 5.95, 5.8)
    add_rect(s, 7.0, 1.22, 5.05, 4.65, "white", "line")
    add_text(s, "核心闭环", 7.25, 1.5, 1.9, 0.35, 23, "teal", True)
    add_bullets(s, ["视频输入：支持摄像头或本地视频文件。", "模型推理：输出边界框、类别和置信度。", "触发报警：前端提示，后端保存报警记录。", "事件留证：保存事件前后视频片段与关键帧。", "人工处置：工单反馈、通知日志、状态更新。"], 7.25, 1.95, 4.45, 2.75, 16)
    add_rect(s, 7.25, 4.95, 4.45, 0.75, "soft", "line")
    add_text(s, "系统把“检测识别”拓展为“识别 + 留证 + 反馈”的工程化流程。", 7.45, 5.12, 4.05, 0.36, 16, "navy", True, PP_ALIGN.CENTER)

    s = prs.slides.add_slide(blank)
    set_bg(s)
    title_bar(s, "2. 研究方法及过程", "详细实现亮点", 6)
    add_rect(s, 0.42, 0.95, 5.7, 5.5, "white", "line")
    add_text(s, "工程实现要点", 0.67, 1.18, 2.5, 0.35, 23, "teal", True)
    add_bullets(s, ["角色权限控制：区分管理员端与普通用户端，前后端双重校验。", "数据留证机制：GridFS 保存报警视频，history 集合记录索引信息。", "日志与通知：记录审计日志、通知日志，支持消息反馈追踪。", "训练与配置：管理员可发起训练、轮询指标并保存系统参数。"], 0.67, 1.7, 5.1, 2.9, 16)
    add_rect(s, 0.67, 5.05, 5.1, 1.05, "soft", "line")
    add_text(s, "技术栈落到项目中后，覆盖数据库、接口、权限、训练日志和通知链路。", 0.9, 5.3, 4.65, 0.42, 16, "navy", True, PP_ALIGN.CENTER)
    add_picture(s, SHOTS / "train-5173.png", 6.8, 1.22, 2.6, 2.05)
    add_picture(s, SHOTS / "logs-5173.png", 9.65, 1.22, 2.6, 2.05)
    add_picture(s, SHOTS / "settings-5173.png", 6.8, 4.0, 2.6, 2.05)
    add_picture(s, SHOTS / "contacts-5173.png", 9.65, 4.0, 2.6, 2.05)

    s = prs.slides.add_slide(blank)
    set_bg(s)
    title_bar(s, "3. 研究结论", "系统完成情况", 7)
    add_rect(s, 0.52, 1.04, 11.8, 5.35, "white", "line")
    add_text(s, "研究结论概括", 0.8, 1.35, 2.5, 0.35, 23, "teal", True)
    add_bullets(s, ["系统已完成实时检测、报警管理、历史回放、模型训练、系统配置等基础模块。", "扩展了角色权限、用户资料、报警工单、通知日志、健康报告、设备心跳等功能。", "实现报警记录、视频片段、关键帧截图、用户消息和日志审计的全流程闭环。", "项目验证表明，该设计能够满足毕业设计场景下的核心功能目标，并具备扩展基础。"], 0.8, 1.85, 6.2, 2.55, 17)
    add_rect(s, 0.8, 4.85, 6.2, 1.05, "navy")
    add_text(s, "结论定位：本课题不只停留在算法验证，而是完成了一个可运行、可管理、可留证的原型系统。", 1.05, 5.12, 5.7, 0.44, 16, "white", True, PP_ALIGN.CENTER)
    add_picture(s, SHOTS / "login-5173.png", 7.7, 1.55, 3.9, 2.2)
    add_picture(s, SHOTS / "statistics-5173.png", 7.7, 4.15, 3.9, 1.85)

    s = prs.slides.add_slide(blank)
    set_bg(s)
    title_bar(s, "3. 研究结论", "测试结果与效果分析", 8)
    add_rect(s, 0.5, 1.0, 5.7, 5.45, "white", "line")
    add_text(s, "测试结果", 0.75, 1.25, 2.0, 0.35, 23, "teal", True)
    add_bullets(s, ["登录与权限测试：管理员与普通用户菜单、接口权限控制均符合预期。", "实时检测测试：可切换视频源，支持阈值调整并能触发报警。", "报警与回放测试：报警记录、留证视频、关键帧截图与删除操作均可用。", "管理功能测试：统计分析、日志中心、模型训练与系统配置协同正常。"], 0.75, 1.75, 5.05, 2.9, 16)
    add_rect(s, 0.75, 5.1, 5.05, 0.9, "soft", "line")
    add_text(s, "测试结论：系统已具备较完整的“输入-识别-显示-报警-留证-处置”能力。", 1.0, 5.33, 4.55, 0.38, 16, "navy", True, PP_ALIGN.CENTER)
    add_picture(s, SHOTS / "statistics-5173.png", 6.85, 1.35, 5.0, 2.65)
    add_picture(s, SHOTS / "logs-5173.png", 6.85, 4.3, 5.0, 1.85)

    s = prs.slides.add_slide(blank)
    set_bg(s)
    title_bar(s, "4. 创新及不足", "项目价值与后续改进方向", 9)
    add_rect(s, 0.6, 1.1, 5.45, 5.0, "white", "line")
    add_rect(s, 6.85, 1.1, 5.45, 5.0, "white", "line")
    add_text(s, "创新点", 0.9, 1.4, 1.8, 0.35, 23, "teal", True)
    add_bullets(s, ["将跌倒检测与报警管理、历史留证、工单反馈整合为统一平台。", "采用前后端分离架构，便于界面扩展、接口维护和业务迭代。", "补充角色权限、通知日志、训练管理等工程化设计。", "围绕养老监护需求，突出“识别 + 留证 + 反馈”的业务闭环。"], 0.9, 1.9, 4.8, 2.9, 16)
    add_text(s, "不足与展望", 7.15, 1.4, 2.2, 0.35, 23, "accent", True)
    add_bullets(s, ["当前密码仍为明文存储，后续需引入哈希加密和更规范的认证机制。", "多设备并发接入、长时间稳定运行与高并发访问能力仍待增强。", "统计分析以基础数量统计为主，风险预测能力有待扩展。", "通知链路依赖第三方服务配置，跨环境部署时需要进一步标准化。"], 7.15, 1.9, 4.8, 2.9, 16)
    add_rect(s, 0.9, 5.25, 10.9, 0.76, "navy")
    add_text(s, "后续可沿模型优化、安全加固、多设备协同和行为趋势预测四个方向继续深化。", 1.2, 5.47, 10.3, 0.32, 18, "white", True, PP_ALIGN.CENTER)

    s = prs.slides.add_slide(blank)
    set_bg(s, "navy")
    add_rect(s, 1.5, 1.25, 10.3, 4.4, "teal")
    add_text(s, "感谢各位老师聆听", 2.5, 2.3, 8.3, 0.7, 34, "white", True, PP_ALIGN.CENTER)
    add_text(s, "恳请批评指正", 2.5, 3.25, 8.3, 0.55, 24, "gold", False, PP_ALIGN.CENTER)
    add_rect(s, 4.9, 4.25, 3.5, 0.05, "accent", radius=False)
    add_text(s, "Q & A", 2.5, 4.65, 8.3, 0.7, 30, "white", True, PP_ALIGN.CENTER)

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
